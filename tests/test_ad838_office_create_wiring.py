"""AD-838 tests — office create-intent wiring, output-path honoring, NL synthesis."""

from __future__ import annotations

import json
from pathlib import Path
from types import SimpleNamespace

import pytest
from docx import Document
from openpyxl import Workbook, load_workbook
from pptx import Presentation

from probos.skill_framework import DocxAgent, PptxAgent, XlsxAgent
from probos.types import IntentMessage


class _SlideLLMClient:
    """Returns a strict-JSON slide array for synthesis."""

    async def complete(self, prompt: str) -> str:
        return json.dumps(
            [
                {"title": "Findings", "bullets": ["Revenue up", "Churn down"]},
                {"title": "Next Steps", "bullets": ["Hire", "Ship"]},
            ]
        )


def _runtime_with_output_dir(output_dir: Path) -> SimpleNamespace:
    office_cfg = SimpleNamespace(output_dir=str(output_dir))
    config = SimpleNamespace(office_skills=office_cfg)
    return SimpleNamespace(config=config)


@pytest.mark.asyncio
async def test_pptx_create_dispatch_produces_readable_deck(tmp_path: Path) -> None:
    agent = PptxAgent(pool="office")
    out = tmp_path / "deck.pptx"
    intent = IntentMessage(
        intent="pptx_create",
        params={
            "title": "Q2 Review",
            "slides": [{"title": "Highlights", "bullets": ["Growth"]}],
            "output_path": str(out),
        },
    )

    result = await agent.handle_intent(intent)

    assert result is not None
    assert result.success is True
    assert result.result["path"] == str(out)
    merged = "\n".join(
        shape.text
        for slide in Presentation(str(out)).slides
        for shape in slide.shapes
        if hasattr(shape, "text")
    )
    assert "Q2 Review" in merged
    assert "Highlights" in merged


@pytest.mark.asyncio
async def test_docx_create_dispatch_produces_readable_doc(tmp_path: Path) -> None:
    agent = DocxAgent(pool="office")
    out = tmp_path / "report.docx"
    intent = IntentMessage(
        intent="docx_create",
        params={
            "title": "Weekly Report",
            "content": ["Item one", "Item two"],
            "output_path": str(out),
        },
    )

    result = await agent.handle_intent(intent)

    assert result is not None and result.success is True
    text = "\n".join(p.text for p in Document(str(out)).paragraphs)
    assert "Weekly Report" in text
    assert "Item one" in text


@pytest.mark.asyncio
async def test_xlsx_update_dispatch_round_trips_cell(tmp_path: Path) -> None:
    workbook_path = tmp_path / "metrics.xlsx"
    wb = Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    ws["A1"] = "Metric"
    ws["B1"] = 1
    wb.save(str(workbook_path))

    agent = XlsxAgent(pool="office")
    intent = IntentMessage(
        intent="xlsx_update",
        params={"file_path": str(workbook_path), "sheet": "Sheet1", "updates": {"B1": 99}},
    )

    result = await agent.handle_intent(intent)

    assert result is not None and result.success is True
    reloaded = load_workbook(result.result["path"])
    assert reloaded["Sheet1"]["B1"].value == 99


@pytest.mark.asyncio
async def test_unknown_intent_declines_returns_none(tmp_path: Path) -> None:
    agent = PptxAgent(pool="office")
    intent = IntentMessage(intent="unrelated_intent", params={})

    result = await agent.handle_intent(intent)

    assert result is None


@pytest.mark.asyncio
async def test_output_path_honored_exact_location(tmp_path: Path) -> None:
    agent = DocxAgent(pool="office")
    out = tmp_path / "nested" / "exact.docx"
    intent = IntentMessage(
        intent="docx_create",
        params={"title": "Exact", "content": ["body"], "output_path": str(out)},
    )

    result = await agent.handle_intent(intent)

    assert result is not None and result.success is True
    assert out.exists()
    assert result.result["path"] == str(out)


@pytest.mark.asyncio
async def test_default_output_dir_slugified_name(tmp_path: Path) -> None:
    output_dir = tmp_path / "office-out"
    agent = DocxAgent(pool="office", runtime=_runtime_with_output_dir(output_dir))
    intent = IntentMessage(
        intent="docx_create",
        params={"title": "My Big Report", "content": ["body"]},
    )

    result = await agent.handle_intent(intent)

    assert result is not None and result.success is True
    expected = output_dir / "my-big-report.docx"
    assert result.result["path"] == str(expected)
    assert expected.exists()


@pytest.mark.asyncio
async def test_nl_synthesis_populates_slides(tmp_path: Path) -> None:
    agent = PptxAgent(pool="office", llm_client=_SlideLLMClient())
    out = tmp_path / "synth.pptx"
    intent = IntentMessage(
        intent="pptx_create",
        params={"title": "Auto Deck", "prompt": "Summarize Q2", "output_path": str(out)},
    )

    result = await agent.handle_intent(intent)

    assert result is not None and result.success is True
    merged = "\n".join(
        shape.text
        for slide in Presentation(str(out)).slides
        for shape in slide.shapes
        if hasattr(shape, "text")
    )
    assert "Findings" in merged
    assert "Revenue up" in merged


@pytest.mark.asyncio
async def test_nl_synthesis_honest_degrade_single_slide(tmp_path: Path) -> None:
    agent = PptxAgent(pool="office", llm_client=None)
    out = tmp_path / "degrade.pptx"
    intent = IntentMessage(
        intent="pptx_create",
        params={"title": "Fallback Deck", "prompt": "no llm available", "output_path": str(out)},
    )

    result = await agent.handle_intent(intent)

    assert result is not None and result.success is True
    deck = Presentation(str(out))
    assert len(deck.slides) == 1


@pytest.mark.asyncio
async def test_missing_required_param_returns_failure(tmp_path: Path) -> None:
    agent = PptxAgent(pool="office")
    intent = IntentMessage(intent="pptx_create", params={"slides": []})

    result = await agent.handle_intent(intent)

    assert result is not None
    assert result.success is False
    assert result.error


def test_consensus_flags_on_mutating_descriptors() -> None:
    def _flag(agent_cls, intent_name: str) -> bool:
        return next(
            d.requires_consensus
            for d in agent_cls.intent_descriptors
            if d.name == intent_name
        )

    assert _flag(DocxAgent, "docx_create") is True
    assert _flag(DocxAgent, "docx_revise") is True
    assert _flag(PptxAgent, "pptx_create") is True
    assert _flag(XlsxAgent, "xlsx_update") is True

    assert _flag(DocxAgent, "docx_summarize") is False
    assert _flag(PptxAgent, "pptx_summarize") is False
    assert _flag(XlsxAgent, "xlsx_read_range") is False
