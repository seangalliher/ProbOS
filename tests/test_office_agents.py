"""AD-755 tests for office document skill agents."""

from __future__ import annotations

from pathlib import Path
from typing import Any

import pytest
from docx import Document
from openpyxl import Workbook, load_workbook
from pptx import Presentation

from probos.skill_framework import DocxAgent, PptxAgent, XlsxAgent


class _FakeLLMClient:
    async def complete(self, prompt: str, **_kwargs: Any) -> str:
        return "- Summary A\n- Summary B\n- Summary C"


@pytest.mark.asyncio
async def test_summarize_docx_extract_text_returns_summary(tmp_path: Path) -> None:
    doc_path = tmp_path / "notes.docx"
    doc = Document()
    doc.add_paragraph("Captain briefing notes")
    doc.add_paragraph("Status is green")
    doc.save(str(doc_path))

    agent = DocxAgent(pool="office")
    summary = await agent.summarize_docx(str(doc_path))

    assert "Summary:" in summary
    assert "Captain briefing notes" in summary


@pytest.mark.asyncio
async def test_summarize_docx_with_llm_client_uses_llm(tmp_path: Path) -> None:
    doc_path = tmp_path / "llm.docx"
    doc = Document()
    doc.add_paragraph("Important text")
    doc.save(str(doc_path))

    agent = DocxAgent(pool="office", llm_client=_FakeLLMClient())
    summary = await agent.summarize_docx(str(doc_path))

    assert summary.startswith("- Summary A")


@pytest.mark.asyncio
async def test_create_docx_from_template_file_created_and_readable(tmp_path: Path) -> None:
    template_path = tmp_path / "template.docx"
    template = Document()
    template.add_paragraph("Template seed")
    template.save(str(template_path))

    agent = DocxAgent(pool="office")
    output_path = await agent.create_docx(
        title="Weekly Report",
        content=["Item one", "Item two"],
        template=str(template_path),
    )

    created = Document(output_path)
    text = "\n".join(p.text for p in created.paragraphs)
    assert "Weekly Report" in text
    assert "Item one" in text


@pytest.mark.asyncio
async def test_summarize_pptx_slide_extraction_works(tmp_path: Path) -> None:
    pptx_path = tmp_path / "briefing.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Mission Update"
    slide.placeholders[1].text = "All systems nominal"
    presentation.save(str(pptx_path))

    agent = PptxAgent(pool="office")
    summary = await agent.summarize_pptx(str(pptx_path))

    assert "PPTX Summary:" in summary
    assert "Mission Update" in summary


@pytest.mark.asyncio
async def test_create_pptx_slide_content_preserved(tmp_path: Path) -> None:
    agent = PptxAgent(pool="office")
    output_path = await agent.create_pptx(
        title="Q2 Review",
        slides=[{"title": "Highlights", "bullets": ["Growth", "Retention"]}],
    )

    created = Presentation(output_path)
    slide_text: list[str] = []
    for slide in created.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text)
    merged = "\n".join(slide_text)
    assert "Q2 Review" in merged
    assert "Highlights" in merged
    assert "Growth" in merged


@pytest.mark.asyncio
async def test_read_update_xlsx_round_trip_integrity(tmp_path: Path) -> None:
    workbook_path = tmp_path / "metrics.xlsx"
    wb = Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    ws["A1"] = "Metric"
    ws["B1"] = "Value"
    ws["A2"] = "Latency"
    ws["B2"] = 42
    wb.save(str(workbook_path))

    agent = XlsxAgent(pool="office")
    before = await agent.read_xlsx_range(str(workbook_path), "Sheet1", "A1:B2")
    assert before == [["Metric", "Value"], ["Latency", 42]]

    updated_path = await agent.update_xlsx(
        str(workbook_path),
        "Sheet1",
        {"B2": 43, "A3": "Errors", "B3": 0},
    )

    after = await agent.read_xlsx_range(updated_path, "Sheet1", "A1:B3")
    assert after == [["Metric", "Value"], ["Latency", 43], ["Errors", 0]]
