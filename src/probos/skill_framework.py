"""AD-428: Agent Skill Framework — developmental competency model.

Data model for skill definitions, proficiency tracking, and agent skill profiles.
Foundation layer — no LLM calls, no I/O in data classes.
"""

from __future__ import annotations

import json
import logging
import os
import re
import tempfile
import time
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Callable

import aiosqlite

from probos.protocols import ConnectionFactory, DatabaseConnection
from probos.substrate.skill_agent import SkillBasedAgent
from probos.tools.protocol import ToolPreference
from probos.types import IntentDescriptor, IntentMessage, IntentResult

logger = logging.getLogger(__name__)


def _slugify(title: str) -> str:
    """Slugify a title into a filesystem-safe base filename (no extension)."""
    slug = re.sub(r"[^a-z0-9]+", "-", (title or "").strip().lower()).strip("-")
    return slug[:80]


def _parse_json_list(raw: Any) -> list[Any] | None:
    """Parse an LLM response into a JSON list, tolerating markdown code fences.

    Returns ``None`` when the response is not a JSON array, so callers can
    honest-degrade to a stub.
    """
    if not isinstance(raw, str):
        return None
    text = raw.strip()
    if text.startswith("```"):
        text = re.sub(r"^```[a-zA-Z0-9]*\n?", "", text)
        text = re.sub(r"\n?```$", "", text).strip()
    try:
        parsed = json.loads(text)
    except (ValueError, TypeError):
        return None
    return parsed if isinstance(parsed, list) else None


def _resolve_default_output_path(runtime: Any, title: str, suffix: str) -> str | None:
    """Resolve a default output path under ``OfficeSkillsConfig.output_dir``.

    Returns ``None`` when no runtime/config output directory is available, so the
    caller falls back to the random-tempfile path (regression-safe for direct
    create_* calls that pass neither ``output_path`` nor a resolvable config).
    """
    config = getattr(runtime, "config", None)
    office_cfg = getattr(config, "office_skills", None)
    output_dir = getattr(office_cfg, "output_dir", None)
    if not output_dir:
        return None
    base = _slugify(title) or "document"
    directory = Path(os.path.expanduser(output_dir))
    return str(directory / f"{base}{suffix}")


class DocxAgent(SkillBasedAgent):
    """DOCX summarization, creation, revision via python-docx."""

    agent_type = "office_docx"
    callsign = "DOCX"
    tier = "domain"
    intent_descriptors = [
        IntentDescriptor(
            name="docx_summarize",
            description="Summarize a DOCX document",
            params={"file_path": "Path to the DOCX file"},
            requires_consensus=False,
        ),
        IntentDescriptor(
            name="docx_create",
            description="Create a DOCX document",
            params={"title": "Document title", "content": "Paragraph list"},
            requires_consensus=True,
        ),
        IntentDescriptor(
            name="docx_revise",
            description="Revise an existing DOCX document",
            params={"file_path": "Path to the DOCX file", "instructions": "Revision instructions"},
            requires_consensus=True,
        ),
    ]

    async def handle_intent(self, intent: IntentMessage) -> IntentResult | None:
        """Dispatch DOCX intents to bound methods; decline unowned intents."""
        params = intent.params or {}
        if intent.intent == "docx_summarize":
            file_path = params.get("file_path")
            if not file_path:
                return IntentResult(
                    intent_id=intent.id, agent_id=self.id, success=False,
                    error="docx_summarize requires 'file_path'",
                )
            summary = await self.summarize_docx(str(file_path))
            return IntentResult(
                intent_id=intent.id, agent_id=self.id, success=True,
                result={"summary": summary},
            )
        if intent.intent == "docx_create":
            title = params.get("title")
            if not title:
                return IntentResult(
                    intent_id=intent.id, agent_id=self.id, success=False,
                    error="docx_create requires 'title'",
                )
            content = params.get("content")
            if not content:
                brief = params.get("prompt") or params.get("brief")
                content = await self._synthesize_paragraphs(str(brief), str(title)) if brief else []
            output_path = params.get("output_path") or _resolve_default_output_path(
                self._runtime, str(title), ".docx"
            )
            path = await self.create_docx(
                title=str(title),
                content=list(content),
                template=params.get("template"),
                output_path=output_path,
            )
            return IntentResult(
                intent_id=intent.id, agent_id=self.id, success=True,
                result={"path": path},
            )
        if intent.intent == "docx_revise":
            file_path = params.get("file_path")
            instructions = params.get("instructions")
            if not file_path or not instructions:
                return IntentResult(
                    intent_id=intent.id, agent_id=self.id, success=False,
                    error="docx_revise requires 'file_path' and 'instructions'",
                )
            path = await self.revise_docx(str(file_path), str(instructions))
            return IntentResult(
                intent_id=intent.id, agent_id=self.id, success=True,
                result={"path": path},
            )
        return None

    async def summarize_docx(self, file_path: str) -> str:
        """Extract full text from .docx and return a concise summary."""
        from docx import Document

        source = Path(file_path)
        if source.suffix.lower() != ".docx":
            raise ValueError("summarize_docx requires a .docx file")
        if not source.exists():
            raise FileNotFoundError(file_path)

        document = Document(str(source))
        chunks: list[str] = [p.text.strip() for p in document.paragraphs if p.text.strip()]
        for table in document.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    chunks.append(row_text)

        full_text = "\n".join(chunks).strip()
        if not full_text:
            return "Document is empty."

        llm_client = getattr(self, "_llm_client", None)
        if llm_client is not None and hasattr(llm_client, "complete"):
            try:
                prompt = (
                    "Summarize this office document in 3 concise bullet points:\n\n"
                    + full_text[:6000]
                )
                maybe_result = llm_client.complete(prompt)
                if hasattr(maybe_result, "__await__"):
                    maybe_result = await maybe_result
                if isinstance(maybe_result, str) and maybe_result.strip():
                    return maybe_result.strip()
            except Exception:
                logger.debug("DocxAgent LLM summarize failed; falling back", exc_info=True)

        return self._summarize_text(full_text)

    async def create_docx(
        self,
        title: str,
        content: list[str],
        template: str | None = None,
        output_path: str | None = None,
    ) -> str:
        """Create new .docx from template or blank document."""
        from docx import Document

        doc = Document(template) if template else Document()
        doc.add_heading(title, level=1)
        for paragraph in content:
            doc.add_paragraph(paragraph)

        if output_path:
            target = Path(os.path.expanduser(output_path))
            target.parent.mkdir(parents=True, exist_ok=True)
            doc.save(str(target))
            return str(target)

        with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as tmp:
            tmp_path = tmp.name
        doc.save(tmp_path)
        return tmp_path

    async def _synthesize_paragraphs(self, brief: str, title: str) -> list[str]:
        """Synthesize document paragraphs from a natural-language brief via the LLM.

        Honest-degrades to a single one-paragraph stub when no LLM client is
        attached or the response cannot be parsed as a JSON string array.
        """
        llm_client = getattr(self, "_llm_client", None)
        if llm_client is not None and hasattr(llm_client, "complete"):
            prompt = (
                "Write the body paragraphs for a document titled "
                f"\"{title}\" based on this brief:\n\n{brief[:4000]}\n\n"
                "Respond with ONLY a strict JSON array of paragraph strings, "
                'e.g. ["First paragraph.", "Second paragraph."]. No prose, no markdown.'
            )
            try:
                maybe = llm_client.complete(prompt)
                if hasattr(maybe, "__await__"):
                    maybe = await maybe
                parsed = _parse_json_list(maybe)
                if parsed is not None:
                    paragraphs = [str(item) for item in parsed if str(item).strip()]
                    if paragraphs:
                        return paragraphs
            except Exception:
                logger.warning(
                    "DocxAgent paragraph synthesis failed; falling back to stub",
                    exc_info=True,
                )
        return [brief.strip() or title]

    async def revise_docx(self, file_path: str, instructions: str) -> str:
        """Apply revision instructions to an existing .docx file."""
        from docx import Document

        source = Path(file_path)
        if source.suffix.lower() != ".docx":
            raise ValueError("revise_docx requires a .docx file")
        if not source.exists():
            raise FileNotFoundError(file_path)

        doc = Document(str(source))
        doc.add_paragraph(f"Revision instructions: {instructions}")
        doc.save(str(source))
        return str(source)

    def _summarize_text(self, text: str) -> str:
        """Return deterministic summary text for extracted document content."""
        snippet = text[:240].replace("\n", " ").strip()
        if len(text) > 240:
            snippet += "..."
        return f"Summary: {snippet}"


class PptxAgent(SkillBasedAgent):
    """PPTX creation and summarization via python-pptx."""

    agent_type = "office_pptx"
    callsign = "PPTX"
    tier = "domain"
    intent_descriptors = [
        IntentDescriptor(
            name="pptx_summarize",
            description="Summarize a PPTX deck",
            params={"file_path": "Path to the PPTX file"},
            requires_consensus=False,
        ),
        IntentDescriptor(
            name="pptx_create",
            description="Create a PPTX deck",
            params={"title": "Deck title", "slides": "Slide definitions"},
            requires_consensus=True,
        ),
    ]

    async def handle_intent(self, intent: IntentMessage) -> IntentResult | None:
        """Dispatch PPTX intents to bound methods; decline unowned intents."""
        params = intent.params or {}
        if intent.intent == "pptx_summarize":
            file_path = params.get("file_path")
            if not file_path:
                return IntentResult(
                    intent_id=intent.id, agent_id=self.id, success=False,
                    error="pptx_summarize requires 'file_path'",
                )
            summary = await self.summarize_pptx(str(file_path))
            return IntentResult(
                intent_id=intent.id, agent_id=self.id, success=True,
                result={"summary": summary},
            )
        if intent.intent == "pptx_create":
            title = params.get("title")
            if not title:
                return IntentResult(
                    intent_id=intent.id, agent_id=self.id, success=False,
                    error="pptx_create requires 'title'",
                )
            slides = params.get("slides")
            if not slides:
                brief = params.get("prompt") or params.get("brief")
                slides = await self._synthesize_slides(str(brief), str(title)) if brief else []
            output_path = params.get("output_path") or _resolve_default_output_path(
                self._runtime, str(title), ".pptx"
            )
            path = await self.create_pptx(
                title=str(title),
                slides=list(slides),
                template=params.get("template"),
                output_path=output_path,
            )
            return IntentResult(
                intent_id=intent.id, agent_id=self.id, success=True,
                result={"path": path},
            )
        return None

    async def summarize_pptx(self, file_path: str) -> str:
        """Extract slide titles and notes from .pptx and return summary text."""
        from pptx import Presentation

        source = Path(file_path)
        if source.suffix.lower() != ".pptx":
            raise ValueError("summarize_pptx requires a .pptx file")
        if not source.exists():
            raise FileNotFoundError(file_path)

        deck = Presentation(str(source))
        lines: list[str] = []
        for idx, slide in enumerate(deck.slides, start=1):
            title_text = ""
            if slide.shapes.title is not None and slide.shapes.title.text:
                title_text = slide.shapes.title.text.strip()
            note_text = ""
            notes_slide = slide.notes_slide
            if notes_slide is not None and notes_slide.notes_text_frame is not None:
                note_text = (notes_slide.notes_text_frame.text or "").strip()
            lines.append(f"Slide {idx}: {title_text}".strip())
            if note_text:
                lines.append(f"Notes: {note_text}")

        body = "\n".join(line for line in lines if line.strip())
        if not body:
            return "Presentation is empty."
        snippet = body[:260]
        if len(body) > 260:
            snippet += "..."
        return f"PPTX Summary: {snippet}"

    async def create_pptx(
        self,
        title: str,
        slides: list[dict[str, Any]],
        template: str | None = None,
        output_path: str | None = None,
    ) -> str:
        """Create a .pptx deck with title and slide content."""
        from pptx import Presentation

        deck = Presentation(template) if template else Presentation()

        title_layout = deck.slide_layouts[0]
        title_slide = deck.slides.add_slide(title_layout)
        if title_slide.shapes.title is not None:
            title_slide.shapes.title.text = title

        for slide_data in slides:
            layout = deck.slide_layouts[1] if len(deck.slide_layouts) > 1 else deck.slide_layouts[0]
            slide = deck.slides.add_slide(layout)
            heading = str(slide_data.get("title", ""))
            body_lines = slide_data.get("bullets") or slide_data.get("content") or []
            if slide.shapes.title is not None:
                slide.shapes.title.text = heading

            placeholder = slide.placeholders[1] if len(slide.placeholders) > 1 else None
            if placeholder is not None and hasattr(placeholder, "text_frame"):
                text_frame = placeholder.text_frame
                text_frame.clear()
                for line in body_lines:
                    paragraph = text_frame.add_paragraph()
                    paragraph.text = str(line)

        if output_path:
            target = Path(os.path.expanduser(output_path))
            target.parent.mkdir(parents=True, exist_ok=True)
            deck.save(str(target))
            return str(target)

        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
            tmp_path = tmp.name
        deck.save(tmp_path)
        return tmp_path

    async def _synthesize_slides(self, brief: str, title: str) -> list[dict[str, Any]]:
        """Synthesize slide definitions from a natural-language brief via the LLM.

        Honest-degrades to an empty slide list (title slide only) when no LLM
        client is attached or the response cannot be parsed.
        """
        llm_client = getattr(self, "_llm_client", None)
        if llm_client is not None and hasattr(llm_client, "complete"):
            prompt = (
                "Design the content slides for a deck titled "
                f"\"{title}\" based on this brief:\n\n{brief[:4000]}\n\n"
                "Respond with ONLY a strict JSON array of slide objects, each "
                '{"title": "Slide title", "bullets": ["point one", "point two"]}. '
                "No prose, no markdown."
            )
            try:
                maybe = llm_client.complete(prompt)
                if hasattr(maybe, "__await__"):
                    maybe = await maybe
                parsed = _parse_json_list(maybe)
                if parsed is not None:
                    slides = [item for item in parsed if isinstance(item, dict)]
                    if slides:
                        return slides
            except Exception:
                logger.warning(
                    "PptxAgent slide synthesis failed; falling back to title slide only",
                    exc_info=True,
                )
        return []


class XlsxAgent(SkillBasedAgent):
    """XLSX operations: read ranges and update values via openpyxl."""

    agent_type = "office_xlsx"
    callsign = "XLSX"
    tier = "domain"
    intent_descriptors = [
        IntentDescriptor(
            name="xlsx_read_range",
            description="Read a range from an XLSX sheet",
            params={"file_path": "Workbook path", "sheet": "Sheet name", "range": "Cell range"},
            requires_consensus=False,
        ),
        IntentDescriptor(
            name="xlsx_update",
            description="Update XLSX cells or formulas",
            params={"file_path": "Workbook path", "sheet": "Sheet name", "updates": "Cell updates"},
            requires_consensus=True,
        ),
    ]

    async def handle_intent(self, intent: IntentMessage) -> IntentResult | None:
        """Dispatch XLSX intents to bound methods; decline unowned intents."""
        params = intent.params or {}
        if intent.intent == "xlsx_read_range":
            file_path = params.get("file_path")
            sheet = params.get("sheet")
            cell_range = params.get("range")
            if not file_path or not sheet or not cell_range:
                return IntentResult(
                    intent_id=intent.id, agent_id=self.id, success=False,
                    error="xlsx_read_range requires 'file_path', 'sheet', and 'range'",
                )
            values = await self.read_xlsx_range(str(file_path), str(sheet), str(cell_range))
            return IntentResult(
                intent_id=intent.id, agent_id=self.id, success=True,
                result={"values": values},
            )
        if intent.intent == "xlsx_update":
            file_path = params.get("file_path")
            sheet = params.get("sheet")
            updates = params.get("updates")
            if not file_path or not sheet or not updates:
                return IntentResult(
                    intent_id=intent.id, agent_id=self.id, success=False,
                    error="xlsx_update requires 'file_path', 'sheet', and 'updates'",
                )
            path = await self.update_xlsx(str(file_path), str(sheet), dict(updates))
            return IntentResult(
                intent_id=intent.id, agent_id=self.id, success=True,
                result={"path": path},
            )
        return None

    async def read_xlsx_range(self, file_path: str, sheet: str, range: str) -> list[list[Any]]:
        """Read a rectangular range from a workbook sheet."""
        from openpyxl import load_workbook

        source = Path(file_path)
        if source.suffix.lower() != ".xlsx":
            raise ValueError("read_xlsx_range requires a .xlsx file")
        if not source.exists():
            raise FileNotFoundError(file_path)

        workbook = load_workbook(str(source))
        if sheet not in workbook.sheetnames:
            raise ValueError(f"Sheet not found: {sheet}")
        worksheet = workbook[sheet]
        rows = worksheet[range]
        return [[cell.value for cell in row] for row in rows]

    async def update_xlsx(self, file_path: str, sheet: str, updates: dict[str, Any]) -> str:
        """Update cell values or formulas in a workbook sheet."""
        from openpyxl import load_workbook

        source = Path(file_path)
        if source.suffix.lower() != ".xlsx":
            raise ValueError("update_xlsx requires a .xlsx file")
        if not source.exists():
            raise FileNotFoundError(file_path)

        workbook = load_workbook(str(source))
        if sheet not in workbook.sheetnames:
            raise ValueError(f"Sheet not found: {sheet}")
        worksheet = workbook[sheet]
        for cell_ref, value in updates.items():
            worksheet[cell_ref] = value
        workbook.save(str(source))
        return str(source)


# ---------------------------------------------------------------------------
# Enums
# ---------------------------------------------------------------------------


from enum import Enum as _Enum


class SkillCategory(_Enum):
    """Three-category skill taxonomy."""
    PCC = "pcc"             # Professional Core Competency — universal crew skills
    ROLE = "role"           # Role/designation skills — department-specific
    ACQUIRED = "acquired"   # Self-developed through experience or mentoring


class ProficiencyLevel(int, _Enum):
    """Seven-level proficiency scale (Dreyfus + Bloom + SFIA unified)."""
    FOLLOW = 1    # Novice: follows explicit procedures
    ASSIST = 2    # Adv. Beginner: recognizes patterns, needs supervision
    APPLY = 3     # Competent: executes independently
    ENABLE = 4    # Competent+: analyzes, decomposes, exercises judgment
    ADVISE = 5    # Proficient: holistic awareness, mentors others
    LEAD = 6      # Expert: innovates, designs new approaches
    SHAPE = 7     # Expert+: sets direction for the domain


# ---------------------------------------------------------------------------
# Data classes
# ---------------------------------------------------------------------------


@dataclass
class SkillDefinition:
    """A skill that agents can acquire and develop."""
    skill_id: str               # e.g., "threat_analysis", "ward_room_communication"
    name: str                   # Human-readable display name
    category: SkillCategory
    description: str = ""
    domain: str = "*"           # "security", "engineering", "*" (universal)
    prerequisites: list[str] = field(default_factory=list)  # skill_ids required at APPLY+
    decay_rate_days: int = 14   # Days idle before proficiency drops one level
    origin: str = "built_in"    # "built_in" (PCC), "role", "acquired", "designed"
    preferred_tools: list[ToolPreference] = field(default_factory=list)
    # AD-428b v1: composite-skill membership. When non-empty, this skill is a
    # composite that fires when the agent has APPLY+ on every constituent.
    composite_skill_ids: list[str] = field(default_factory=list)
    # AD-428b v1: pairwise synergy declaration. When skill A lists B AND B lists A,
    # SkillProfile.synergy_bonus(A, B) returns a non-zero float.
    synergy_partners: list[str] = field(default_factory=list)


@dataclass
class AgentSkillRecord:
    """An agent's proficiency in a specific skill."""
    agent_id: str
    skill_id: str
    proficiency: ProficiencyLevel = ProficiencyLevel.FOLLOW
    acquired_at: float = 0.0
    last_exercised: float = 0.0
    exercise_count: int = 0
    acquisition_source: str = "commissioning"  # "commissioning", "qualification", "experience", "mentoring"
    suspended: bool = False     # True if model lacks required capabilities
    assessment_history: list[dict] = field(default_factory=list)
        # [{timestamp, level, source, notes}]

    def to_dict(self) -> dict[str, Any]:
        """Serialize for API/persistence."""
        return {
            "agent_id": self.agent_id,
            "skill_id": self.skill_id,
            "proficiency": self.proficiency.value,
            "proficiency_label": self.proficiency.name.lower(),
            "acquired_at": self.acquired_at,
            "last_exercised": self.last_exercised,
            "exercise_count": self.exercise_count,
            "acquisition_source": self.acquisition_source,
            "suspended": self.suspended,
        }


@dataclass
class SkillProfile:
    """Complete skill profile for an agent."""
    agent_id: str
    pccs: list[AgentSkillRecord] = field(default_factory=list)
    role_skills: list[AgentSkillRecord] = field(default_factory=list)
    acquired_skills: list[AgentSkillRecord] = field(default_factory=list)

    @property
    def all_skills(self) -> list[AgentSkillRecord]:
        return self.pccs + self.role_skills + self.acquired_skills

    @property
    def depth(self) -> int:
        """Max proficiency across all skills."""
        if not self.all_skills:
            return 0
        return max(s.proficiency.value for s in self.all_skills)

    @property
    def breadth(self) -> int:
        """Number of distinct domains with ASSIST+ proficiency."""
        domains = set()
        for s in self.all_skills:
            if s.proficiency.value >= ProficiencyLevel.ASSIST.value and not s.suspended:
                domains.add(s.skill_id)
        return len(domains)

    def get_proficiency(self, skill_id: str) -> "ProficiencyLevel | None":
        """AD-428b v1: lookup proficiency for a skill_id; None if not held or suspended."""
        for record in self.all_skills:
            if record.skill_id == skill_id and not record.suspended:
                return record.proficiency
        return None

    def has_composite_capability(
        self, composite: "SkillDefinition"
    ) -> bool:
        """AD-428b v1: True iff agent has APPLY+ on EVERY constituent of the composite.

        Composites with empty composite_skill_ids never fire (degenerate case).
        """
        if not composite.composite_skill_ids:
            return False
        for constituent_id in composite.composite_skill_ids:
            level = self.get_proficiency(constituent_id)
            if level is None or level.value < ProficiencyLevel.APPLY.value:
                return False
        return True

    def synergy_bonus(
        self,
        skill_a_id: str,
        skill_b_id: str,
        registry_lookup: Any = None,
    ) -> float:
        """AD-428b v1: pairwise synergy bonus between two skills.

        Returns 0.0 unless ALL of:
          - agent holds both skills at APPLY+
          - skill A's SkillDefinition.synergy_partners contains B
          - skill B's SkillDefinition.synergy_partners contains A
        Bonus = 0.10 * (min(level_a, level_b) - APPLY + 1) capped at 0.50.

        registry_lookup: callable taking skill_id -> SkillDefinition | None.
        Pass None to opt out of synergy_partners check (returns 0.0).
        """
        if registry_lookup is None or skill_a_id == skill_b_id:
            return 0.0
        level_a = self.get_proficiency(skill_a_id)
        level_b = self.get_proficiency(skill_b_id)
        if level_a is None or level_b is None:
            return 0.0
        if level_a.value < ProficiencyLevel.APPLY.value or level_b.value < ProficiencyLevel.APPLY.value:
            return 0.0
        defn_a = registry_lookup(skill_a_id)
        defn_b = registry_lookup(skill_b_id)
        if defn_a is None or defn_b is None:
            return 0.0
        if skill_b_id not in defn_a.synergy_partners:
            return 0.0
        if skill_a_id not in defn_b.synergy_partners:
            return 0.0
        bonus = 0.10 * (min(level_a.value, level_b.value) - ProficiencyLevel.APPLY.value + 1)
        return min(bonus, 0.50)

    def to_dict(self) -> dict[str, Any]:
        return {
            "agent_id": self.agent_id,
            "pccs": [s.to_dict() for s in self.pccs],
            "role_skills": [s.to_dict() for s in self.role_skills],
            "acquired_skills": [s.to_dict() for s in self.acquired_skills],
            "depth": self.depth,
            "breadth": self.breadth,
        }


@dataclass
class QualificationRecord:
    """Tracks an agent's progress through a qualification path (AD-429b)."""
    agent_id: str
    path_id: str  # e.g., "ensign_to_lieutenant"
    started_at: float
    completed_at: float | None = None
    requirement_status: dict[str, bool] = field(default_factory=dict)

    def is_complete(self) -> bool:
        return all(self.requirement_status.values()) if self.requirement_status else False

    def to_dict(self) -> dict:
        return {
            "agent_id": self.agent_id,
            "path_id": self.path_id,
            "started_at": self.started_at,
            "completed_at": self.completed_at,
            "is_complete": self.is_complete(),
            "requirements": self.requirement_status,
        }


# ---------------------------------------------------------------------------
# Built-in skill definitions
# ---------------------------------------------------------------------------

# ── Built-in Professional Core Competencies ────────────────────────
BUILTIN_PCCS: list[SkillDefinition] = [
    SkillDefinition(
        skill_id="communication",
        name="Communication",
        category=SkillCategory.PCC,
        description="Effective Ward Room participation, report structure, endorsement quality.",
        domain="*",
        decay_rate_days=30,
    ),
    SkillDefinition(
        skill_id="chain_of_command",
        name="Chain of Command",
        category=SkillCategory.PCC,
        description="Standing Orders compliance, escalation protocols, rank-appropriate behavior.",
        domain="*",
        decay_rate_days=30,
    ),
    SkillDefinition(
        skill_id="duty_execution",
        name="Duty Execution",
        category=SkillCategory.PCC,
        description="Completing scheduled duties on time, structured reporting, prioritization.",
        domain="*",
        decay_rate_days=30,
    ),
    SkillDefinition(
        skill_id="collaboration",
        name="Collaboration",
        category=SkillCategory.PCC,
        description="Consensus participation, cross-agent coordination, constructive disagreement.",
        domain="*",
        decay_rate_days=30,
    ),
    SkillDefinition(
        skill_id="knowledge_stewardship",
        name="Knowledge Stewardship",
        category=SkillCategory.PCC,
        description="Contributing to shared knowledge, accurate episodic recording.",
        domain="*",
        decay_rate_days=30,
    ),
    SkillDefinition(
        skill_id="self_assessment",
        name="Self-Assessment",
        category=SkillCategory.PCC,
        description="Recognizing own limitations, requesting assistance appropriately.",
        domain="*",
        decay_rate_days=30,
    ),
    SkillDefinition(
        skill_id="ethical_reasoning",
        name="Ethical Reasoning",
        category=SkillCategory.PCC,
        description="Standing Orders internalization, safety awareness, reversibility preference.",
        domain="*",
        decay_rate_days=30,
    ),
    # BF-294 / AD-728d: bind the cognitive skill at config/skills/self-image-awareness/SKILL.md
    # to the proficiency-tracking SkillRegistry so AD-596c's skill_bridge sync matches it
    # ("matched" instead of "unmatched"). Without this entry the skill loads in
    # augmentation mode but proficiency gating is inactive — observable as the boot
    # warning "Cognitive skill 'self-image-awareness' references skill_id
    # 'self-image-awareness' not found in SkillRegistry — proficiency gating will be inactive".
    SkillDefinition(
        skill_id="self-image-awareness",
        name="Self-Image Awareness",
        category=SkillCategory.PCC,
        description=(
            "Vision-based avatar self-check via the [SELF_CHECK reason] marker — "
            "invoking the AD-728c render-coherence mirror to verify the rendered "
            "avatar matches declared intent. Universal crew skill; budget-aware "
            "(3/hr idle OR 2/active-conversation INSTEAD OF hourly)."
        ),
        domain="*",
        decay_rate_days=30,
    ),
    # BF-295 / AD-634: same pattern — bind the cognitive skill at
    # config/skills/notebook-quality/SKILL.md to the proficiency-tracking
    # SkillRegistry. AD-634 (Wave 70) shipped the SKILL.md but never added the
    # registry entry; warning has been silently emitting since then.
    SkillDefinition(
        skill_id="notebook-quality",
        name="Notebook Analytical Quality",
        category=SkillCategory.PCC,
        description=(
            "Analytical quality discipline for notebook entries — Analytical "
            "Purpose Gate, Finding-First (Minto Pyramid), Temporal Threading, "
            "Data vs Analysis, Ward Room Differentiation, Pre-Write Verification "
            "Gate. Co-activates with communication-discipline on proactive_think."
        ),
        domain="*",
        decay_rate_days=30,
    ),
]

# ── Role skill templates per department ────────────────────────────
# Keyed by agent_type → list of SkillDefinition
ROLE_SKILL_TEMPLATES: dict[str, list[SkillDefinition]] = {
    "security_officer": [
        SkillDefinition(skill_id="threat_analysis", name="Threat Analysis", category=SkillCategory.ROLE, domain="security", decay_rate_days=14, origin="role"),
        SkillDefinition(skill_id="vulnerability_assessment", name="Vulnerability Assessment", category=SkillCategory.ROLE, domain="security", prerequisites=["threat_analysis"], decay_rate_days=14, origin="role"),
        SkillDefinition(skill_id="audit_procedures", name="Audit Procedures", category=SkillCategory.ROLE, domain="security", decay_rate_days=14, origin="role"),
        SkillDefinition(skill_id="agentic_security_review", name="Agentic Security Review", category=SkillCategory.ROLE, domain="security", prerequisites=["threat_analysis"], decay_rate_days=14, origin="role"),
    ],
    "engineering_officer": [
        SkillDefinition(skill_id="code_review", name="Code Review", category=SkillCategory.ROLE, domain="engineering", decay_rate_days=14, origin="role"),
        SkillDefinition(skill_id="architecture_analysis", name="Architecture Analysis", category=SkillCategory.ROLE, domain="engineering", prerequisites=["code_review"], decay_rate_days=14, origin="role"),
        SkillDefinition(skill_id="performance_optimization", name="Performance Optimization", category=SkillCategory.ROLE, domain="engineering", prerequisites=["architecture_analysis"], decay_rate_days=14, origin="role"),
    ],
    "operations_officer": [
        SkillDefinition(skill_id="resource_management", name="Resource Management", category=SkillCategory.ROLE, domain="operations", decay_rate_days=14, origin="role"),
        SkillDefinition(skill_id="scheduling_optimization", name="Scheduling Optimization", category=SkillCategory.ROLE, domain="operations", decay_rate_days=14, origin="role"),
        SkillDefinition(skill_id="incident_response", name="Incident Response", category=SkillCategory.ROLE, domain="operations", prerequisites=["resource_management"], decay_rate_days=14, origin="role"),
    ],
    "diagnostician": [
        SkillDefinition(skill_id="health_assessment", name="Health Assessment", category=SkillCategory.ROLE, domain="medical", decay_rate_days=14, origin="role"),
        SkillDefinition(skill_id="anomaly_detection", name="Anomaly Detection", category=SkillCategory.ROLE, domain="medical", decay_rate_days=14, origin="role"),
        SkillDefinition(skill_id="diagnostic_reasoning", name="Diagnostic Reasoning", category=SkillCategory.ROLE, domain="medical", prerequisites=["health_assessment", "anomaly_detection"], decay_rate_days=14, origin="role"),
    ],
    "scout": [
        SkillDefinition(skill_id="codebase_exploration", name="Codebase Exploration", category=SkillCategory.ROLE, domain="science", decay_rate_days=7, origin="role"),
        SkillDefinition(skill_id="information_gathering", name="Information Gathering", category=SkillCategory.ROLE, domain="science", decay_rate_days=7, origin="role"),
        SkillDefinition(skill_id="pattern_identification", name="Pattern Identification", category=SkillCategory.ROLE, domain="science", prerequisites=["codebase_exploration"], decay_rate_days=7, origin="role"),
    ],
    "counselor": [
        SkillDefinition(skill_id="cognitive_health_eval", name="Cognitive Health Evaluation", category=SkillCategory.ROLE, domain="medical", decay_rate_days=14, origin="role"),
        SkillDefinition(skill_id="crew_fitness_assessment", name="Crew Fitness Assessment", category=SkillCategory.ROLE, domain="medical", decay_rate_days=14, origin="role"),
        SkillDefinition(skill_id="conflict_mediation", name="Conflict Mediation", category=SkillCategory.ROLE, domain="medical", prerequisites=["cognitive_health_eval"], decay_rate_days=14, origin="role"),
        SkillDefinition(skill_id="graduated_coaching", name="Graduated Coaching", category=SkillCategory.ROLE, domain="medical", prerequisites=["cognitive_health_eval"], decay_rate_days=14, origin="role"),
    ],
    "architect": [
        SkillDefinition(skill_id="design_review", name="Design Review", category=SkillCategory.ROLE, domain="science", decay_rate_days=14, origin="role"),
        SkillDefinition(skill_id="strategic_planning", name="Strategic Planning", category=SkillCategory.ROLE, domain="science", prerequisites=["design_review"], decay_rate_days=14, origin="role"),
        SkillDefinition(skill_id="technology_evaluation", name="Technology Evaluation", category=SkillCategory.ROLE, domain="science", decay_rate_days=14, origin="role"),
    ],
    "builder": [
        SkillDefinition(skill_id="component_integration", name="Component Integration", category=SkillCategory.ROLE, domain="engineering", decay_rate_days=14, origin="role"),
        SkillDefinition(skill_id="build_automation", name="Build Automation", category=SkillCategory.ROLE, domain="engineering", decay_rate_days=14, origin="role"),
        SkillDefinition(skill_id="code_generation", name="Code Generation", category=SkillCategory.ROLE, domain="engineering", prerequisites=["component_integration"], decay_rate_days=14, origin="role"),
    ],
    "surgeon": [
        SkillDefinition(skill_id="surgical_precision", name="Surgical Precision", category=SkillCategory.ROLE, domain="medical", decay_rate_days=14, origin="role"),
        SkillDefinition(skill_id="crisis_response", name="Crisis Response", category=SkillCategory.ROLE, domain="medical", decay_rate_days=14, origin="role"),
        SkillDefinition(skill_id="system_repair", name="System Repair", category=SkillCategory.ROLE, domain="medical", prerequisites=["surgical_precision", "crisis_response"], decay_rate_days=14, origin="role"),
    ],
    "pharmacist": [
        SkillDefinition(skill_id="intervention_management", name="Intervention Management", category=SkillCategory.ROLE, domain="medical", decay_rate_days=14, origin="role"),
        SkillDefinition(skill_id="interaction_analysis", name="Interaction Analysis", category=SkillCategory.ROLE, domain="medical", decay_rate_days=14, origin="role"),
        SkillDefinition(skill_id="compliance_review", name="Compliance Review", category=SkillCategory.ROLE, domain="medical", prerequisites=["intervention_management"], decay_rate_days=14, origin="role"),
    ],
    "pathologist": [
        SkillDefinition(skill_id="system_analysis", name="System Analysis", category=SkillCategory.ROLE, domain="medical", decay_rate_days=14, origin="role"),
        SkillDefinition(skill_id="failure_identification", name="Failure Identification", category=SkillCategory.ROLE, domain="medical", decay_rate_days=14, origin="role"),
        SkillDefinition(skill_id="research_methodology", name="Research Methodology", category=SkillCategory.ROLE, domain="medical", prerequisites=["system_analysis"], decay_rate_days=14, origin="role"),
    ],
    "data_analyst": [
        SkillDefinition(skill_id="data_visualization", name="Data Visualization", category=SkillCategory.ROLE, domain="science", decay_rate_days=7, origin="role"),
        SkillDefinition(skill_id="statistical_analysis", name="Statistical Analysis", category=SkillCategory.ROLE, domain="science", decay_rate_days=7, origin="role"),
        SkillDefinition(skill_id="trend_identification", name="Trend Identification", category=SkillCategory.ROLE, domain="science", prerequisites=["statistical_analysis"], decay_rate_days=7, origin="role"),
    ],
    "systems_analyst": [
        SkillDefinition(skill_id="requirements_analysis", name="Requirements Analysis", category=SkillCategory.ROLE, domain="science", decay_rate_days=7, origin="role"),
        SkillDefinition(skill_id="process_optimization", name="Process Optimization", category=SkillCategory.ROLE, domain="science", decay_rate_days=7, origin="role"),
        SkillDefinition(skill_id="integration_testing", name="Integration Testing", category=SkillCategory.ROLE, domain="science", prerequisites=["requirements_analysis"], decay_rate_days=7, origin="role"),
    ],
    "research_specialist": [
        SkillDefinition(skill_id="literature_review", name="Literature Review", category=SkillCategory.ROLE, domain="science", decay_rate_days=7, origin="role"),
        SkillDefinition(skill_id="hypothesis_testing", name="Hypothesis Testing", category=SkillCategory.ROLE, domain="science", decay_rate_days=7, origin="role"),
        SkillDefinition(skill_id="experimental_design", name="Experimental Design", category=SkillCategory.ROLE, domain="science", prerequisites=["hypothesis_testing"], decay_rate_days=7, origin="role"),
    ],
}


# ---------------------------------------------------------------------------
# SQLite schema
# ---------------------------------------------------------------------------

_SCHEMA = """
CREATE TABLE IF NOT EXISTS skill_definitions (
    skill_id TEXT PRIMARY KEY,
    name TEXT NOT NULL,
    category TEXT NOT NULL,
    description TEXT DEFAULT '',
    domain TEXT DEFAULT '*',
    prerequisites TEXT DEFAULT '[]',
    decay_rate_days INTEGER DEFAULT 14,
    origin TEXT DEFAULT 'built_in'
);

CREATE TABLE IF NOT EXISTS agent_skills (
    agent_id TEXT NOT NULL,
    skill_id TEXT NOT NULL,
    proficiency INTEGER DEFAULT 1,
    acquired_at REAL DEFAULT 0,
    last_exercised REAL DEFAULT 0,
    exercise_count INTEGER DEFAULT 0,
    acquisition_source TEXT DEFAULT 'commissioning',
    suspended INTEGER DEFAULT 0,
    assessment_history TEXT DEFAULT '[]',
    PRIMARY KEY (agent_id, skill_id)
);

CREATE INDEX IF NOT EXISTS idx_agent_skills_agent ON agent_skills(agent_id);
CREATE INDEX IF NOT EXISTS idx_agent_skills_skill ON agent_skills(skill_id);
CREATE INDEX IF NOT EXISTS idx_skill_defs_category ON skill_definitions(category);
CREATE INDEX IF NOT EXISTS idx_skill_defs_domain ON skill_definitions(domain);

CREATE TABLE IF NOT EXISTS qualification_records (
    agent_id TEXT NOT NULL,
    path_id TEXT NOT NULL,
    started_at REAL NOT NULL,
    completed_at REAL,
    requirement_status TEXT NOT NULL DEFAULT '{}',
    PRIMARY KEY (agent_id, path_id)
);

-- AD-428b v1: per-agent development goals (one goal per skill).
CREATE TABLE IF NOT EXISTS agent_development_goals (
    agent_id TEXT NOT NULL,
    skill_id TEXT NOT NULL,
    target_level INTEGER NOT NULL,
    set_at REAL NOT NULL,
    notes TEXT DEFAULT '',
    PRIMARY KEY (agent_id, skill_id)
);

CREATE INDEX IF NOT EXISTS idx_dev_goals_agent ON agent_development_goals(agent_id);
"""


# ---------------------------------------------------------------------------
# SkillRegistry — master catalog of skill definitions
# ---------------------------------------------------------------------------


class SkillRegistry:
    """Ship's Computer service — manages the master catalog of skill definitions.

    Infrastructure tier (no identity). Provides CRUD for SkillDefinitions
    and prerequisite DAG queries.
    """

    def __init__(self, db_path: str | None = None, connection_factory: ConnectionFactory | None = None):
        self._db_path = db_path
        self._db: DatabaseConnection | None = None
        # In-memory cache for fast lookup
        self._cache: dict[str, SkillDefinition] = {}
        self._connection_factory = connection_factory
        if self._connection_factory is None:
            from probos.storage.sqlite_factory import default_factory
            self._connection_factory = default_factory

    async def start(self) -> None:
        if self._db_path:
            self._db = await self._connection_factory.connect(self._db_path)
            await self._db.execute("PRAGMA foreign_keys = ON")
            self._db.row_factory = aiosqlite.Row
            await self._db.executescript(_SCHEMA)
            await self._db.commit()
            # AD-423a: Add preferred_tools column if missing (migration)
            try:
                await self._db.execute(
                    "ALTER TABLE skill_definitions ADD COLUMN preferred_tools TEXT DEFAULT '[]'"
                )
                await self._db.commit()
            except Exception:
                pass  # Column already exists
            # AD-428b v1: composite_skill_ids column (JSON-encoded list)
            try:
                await self._db.execute(
                    "ALTER TABLE skill_definitions ADD COLUMN composite_skill_ids TEXT DEFAULT '[]'"
                )
                await self._db.commit()
            except Exception:
                pass  # Column already exists
            # AD-428b v1: synergy_partners column (JSON-encoded list)
            try:
                await self._db.execute(
                    "ALTER TABLE skill_definitions ADD COLUMN synergy_partners TEXT DEFAULT '[]'"
                )
                await self._db.commit()
            except Exception:
                pass  # Column already exists
            # Load existing definitions into cache
            async with self._db.execute("SELECT * FROM skill_definitions") as cur:
                async for row in cur:
                    self._cache[row["skill_id"]] = self._row_to_definition(row)

    async def stop(self) -> None:
        if self._db:
            await self._db.close()
            self._db = None

    def _row_to_definition(self, row) -> SkillDefinition:
        prefs_raw = json.loads(row["preferred_tools"] if "preferred_tools" in row.keys() else "[]")
        prefs = [ToolPreference(tool_id=p["tool_id"], priority=p.get("priority", 0), context=p.get("context", "")) for p in prefs_raw]
        # AD-428b v1: tolerate older rows where columns are absent.
        composite_raw = row["composite_skill_ids"] if "composite_skill_ids" in row.keys() else "[]"
        synergy_raw = row["synergy_partners"] if "synergy_partners" in row.keys() else "[]"
        return SkillDefinition(
            skill_id=row["skill_id"],
            name=row["name"],
            category=SkillCategory(row["category"]),
            description=row["description"] or "",
            domain=row["domain"] or "*",
            prerequisites=json.loads(row["prerequisites"] or "[]"),
            decay_rate_days=row["decay_rate_days"] or 14,
            origin=row["origin"] or "built_in",
            preferred_tools=prefs,
            composite_skill_ids=json.loads(composite_raw or "[]"),
            synergy_partners=json.loads(synergy_raw or "[]"),
        )

    async def register_from_manifest(self, yaml_path: "Path") -> SkillDefinition:
        """AD-481e: load skill.yaml + register the resulting SkillDefinition.

        Thin composition helper — equivalent to:
            defn = load_skill_from_manifest(yaml_path)
            return await self.register_skill(defn)
        """
        from probos.extensions.skill_manifest import load_skill_from_manifest
        defn = load_skill_from_manifest(yaml_path)
        return await self.register_skill(defn)

    async def register_skill(self, defn: SkillDefinition) -> SkillDefinition:
        """Register or update a skill definition."""
        self._cache[defn.skill_id] = defn
        if self._db:
            prefs_json = json.dumps([{"tool_id": p.tool_id, "priority": p.priority, "context": p.context} for p in defn.preferred_tools])
            composite_json = json.dumps(defn.composite_skill_ids)
            synergy_json = json.dumps(defn.synergy_partners)
            await self._db.execute(
                "INSERT OR REPLACE INTO skill_definitions "
                "(skill_id, name, category, description, domain, prerequisites, decay_rate_days, origin, preferred_tools, composite_skill_ids, synergy_partners) "
                "VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)",
                (defn.skill_id, defn.name, defn.category.value, defn.description,
                 defn.domain, json.dumps(defn.prerequisites), defn.decay_rate_days, defn.origin, prefs_json,
                 composite_json, synergy_json),
            )
            await self._db.commit()
        return defn

    async def register_builtins(self) -> None:
        """Register all built-in PCCs and role skill templates."""
        for pcc in BUILTIN_PCCS:
            await self.register_skill(pcc)
        for role_skills in ROLE_SKILL_TEMPLATES.values():
            for skill in role_skills:
                await self.register_skill(skill)

    def get_skill(self, skill_id: str) -> SkillDefinition | None:
        """Get a skill definition by ID (from cache)."""
        return self._cache.get(skill_id)

    def list_skills(
        self, category: SkillCategory | None = None, domain: str | None = None,
    ) -> list[SkillDefinition]:
        """List skill definitions with optional filters."""
        result = list(self._cache.values())
        if category:
            result = [s for s in result if s.category == category]
        if domain:
            result = [s for s in result if s.domain == domain or s.domain == "*"]
        return sorted(result, key=lambda s: s.skill_id)

    def get_prerequisites(self, skill_id: str) -> list[str]:
        """Get the full prerequisite DAG for a skill (flattened, deduplicated)."""
        visited: set[str] = set()
        result: list[str] = []

        def _walk(sid: str) -> None:
            defn = self._cache.get(sid)
            if not defn:
                return
            for prereq_id in defn.prerequisites:
                if prereq_id not in visited:
                    visited.add(prereq_id)
                    _walk(prereq_id)
                    result.append(prereq_id)

        _walk(skill_id)
        return result


# ---------------------------------------------------------------------------
# AgentSkillService — per-agent skill records
# ---------------------------------------------------------------------------


class AgentSkillService:
    """Ship's Computer service — manages per-agent skill records.

    Infrastructure tier. Tracks acquisition, proficiency, decay, and
    produces SkillProfiles.
    """

    def __init__(self, db_path: str | None = None, registry: SkillRegistry | None = None, connection_factory: ConnectionFactory | None = None):
        self._db_path = db_path
        self._db: DatabaseConnection | None = None
        self._registry = registry
        self._connection_factory = connection_factory
        if self._connection_factory is None:
            from probos.storage.sqlite_factory import default_factory
            self._connection_factory = default_factory
        # AD-628a: skill telemetry emitter slot (registered via set_event_emitter)
        self._event_emitter: Callable[[Any, dict[str, Any]], None] | None = None

    def set_event_emitter(self, emitter: "Callable[[Any, dict[str, Any]], None] | None") -> None:
        """Register an event emitter for skill telemetry (AD-628a).

        The emitter is called with ``(EventType, payload_dict)`` on each
        mutation. Tier-2 log-and-degrade — emitter exceptions are caught
        and logged at warning level. Pass ``None`` to unregister.
        """
        self._event_emitter = emitter

    def _emit(self, event_type: Any, payload: dict[str, Any]) -> None:
        """AD-628a: tier-2 log-and-degrade emission helper."""
        if self._event_emitter is None:
            return
        try:
            self._event_emitter(event_type, payload)
        except Exception:
            logger.warning(
                "AD-628a: skill telemetry emit failed for %s",
                event_type,
                exc_info=True,
            )

    async def start(self) -> None:
        if self._db_path:
            self._db = await self._connection_factory.connect(self._db_path)
            await self._db.execute("PRAGMA foreign_keys = ON")
            self._db.row_factory = aiosqlite.Row
            await self._db.executescript(_SCHEMA)
            await self._db.commit()

    async def stop(self) -> None:
        if self._db:
            await self._db.close()
            self._db = None

    def _row_to_record(self, row) -> AgentSkillRecord:
        return AgentSkillRecord(
            agent_id=row["agent_id"],
            skill_id=row["skill_id"],
            proficiency=ProficiencyLevel(row["proficiency"]),
            acquired_at=row["acquired_at"] or 0.0,
            last_exercised=row["last_exercised"] or 0.0,
            exercise_count=row["exercise_count"] or 0,
            acquisition_source=row["acquisition_source"] or "commissioning",
            suspended=bool(row["suspended"]),
            assessment_history=json.loads(row["assessment_history"] or "[]"),
        )

    async def acquire_skill(
        self,
        agent_id: str,
        skill_id: str,
        source: str = "commissioning",
        proficiency: ProficiencyLevel = ProficiencyLevel.FOLLOW,
    ) -> AgentSkillRecord:
        """Give an agent a skill at a starting proficiency level.

        Raises ValueError if prerequisites are not met.
        """
        # Check prerequisites
        if self._registry:
            defn = self._registry.get_skill(skill_id)
            if defn and defn.prerequisites:
                for prereq_id in defn.prerequisites:
                    existing = await self._get_record(agent_id, prereq_id)
                    if not existing or existing.proficiency.value < ProficiencyLevel.APPLY.value:
                        raise ValueError(
                            f"Prerequisite '{prereq_id}' not met for '{skill_id}' "
                            f"(requires APPLY+, agent has "
                            f"{'none' if not existing else existing.proficiency.name})"
                        )

        now = time.time()
        record = AgentSkillRecord(
            agent_id=agent_id,
            skill_id=skill_id,
            proficiency=proficiency,
            acquired_at=now,
            last_exercised=now,
            exercise_count=0,
            acquisition_source=source,
        )
        await self._upsert_record(record)
        # AD-628a: SKILL_ACQUIRED telemetry
        from probos.events import EventType as _ET
        self._emit(_ET.SKILL_ACQUIRED, {
            "agent_id": agent_id,
            "skill_id": skill_id,
            "to_level": record.proficiency.value,
            "source": source,
            "timestamp": now,
            "reason": "acquired",
        })
        return record

    async def commission_agent(self, agent_id: str, agent_type: str) -> SkillProfile:
        """Assign an agent their initial skill complement (PCCs + role skills).

        Called during agent registration/onboarding.
        """
        # All crew get PCCs at FOLLOW
        for pcc in BUILTIN_PCCS:
            try:
                await self.acquire_skill(
                    agent_id, pcc.skill_id, source="commissioning",
                    proficiency=ProficiencyLevel.FOLLOW,
                )
            except ValueError:
                pass  # Already exists or prereq issue — skip

        # Role-specific skills
        role_skills = ROLE_SKILL_TEMPLATES.get(agent_type, [])
        for skill in role_skills:
            try:
                await self.acquire_skill(
                    agent_id, skill.skill_id, source="commissioning",
                    proficiency=ProficiencyLevel.FOLLOW,
                )
            except ValueError:
                pass  # Prerequisite not met at commissioning — expected for chained skills

        return await self.get_profile(agent_id)

    async def update_proficiency(
        self,
        agent_id: str,
        skill_id: str,
        new_level: ProficiencyLevel,
        source: str = "assessment",
        notes: str = "",
    ) -> AgentSkillRecord | None:
        """Update an agent's proficiency level with assessment record."""
        record = await self._get_record(agent_id, skill_id)
        if not record:
            return None
        from_level = record.proficiency
        to_level = new_level
        record.proficiency = new_level
        record.last_exercised = time.time()
        record.assessment_history.append({
            "timestamp": time.time(),
            "level": new_level.value,
            "source": source,
            "notes": notes,
        })
        await self._upsert_record(record)
        # AD-628a: SKILL_REGRESSION on downward, SKILL_EXERCISED otherwise
        from probos.events import EventType as _ET
        _evt = _ET.SKILL_REGRESSION if to_level.value < from_level.value else _ET.SKILL_EXERCISED
        self._emit(_evt, {
            "agent_id": agent_id,
            "skill_id": skill_id,
            "from_level": from_level.value,
            "to_level": to_level.value,
            "source": source,
            "timestamp": time.time(),
            "reason": notes or "proficiency_update",
        })
        return record

    async def record_exercise(self, agent_id: str, skill_id: str) -> AgentSkillRecord | None:
        """Record that an agent exercised a skill (resets decay timer)."""
        record = await self._get_record(agent_id, skill_id)
        if not record:
            return None
        record.last_exercised = time.time()
        record.exercise_count += 1
        await self._upsert_record(record)
        # AD-628a: SKILL_EXERCISED telemetry
        from probos.events import EventType as _ET
        self._emit(_ET.SKILL_EXERCISED, {
            "agent_id": agent_id,
            "skill_id": skill_id,
            "to_level": record.proficiency.value,
            "source": "exercise",
            "timestamp": record.last_exercised,
            "reason": "recorded_exercise",
        })
        return record

    async def check_decay(self, now: float | None = None) -> list[AgentSkillRecord]:
        """Find all skills that have decayed due to inactivity.

        Returns list of records that were downgraded.
        """
        if now is None:
            now = time.time()
        decayed: list[AgentSkillRecord] = []
        if not self._db:
            return decayed

        async with self._db.execute(
            "SELECT * FROM agent_skills WHERE proficiency > 1 AND suspended = 0"
        ) as cur:
            async for row in cur:
                record = self._row_to_record(row)
                defn = self._registry.get_skill(record.skill_id) if self._registry else None
                decay_days = defn.decay_rate_days if defn else 14
                idle_seconds = now - record.last_exercised
                idle_days = idle_seconds / 86400.0
                if idle_days >= decay_days:
                    # Drop one level per decay period elapsed
                    levels_dropped = int(idle_days / decay_days)
                    new_level = max(1, record.proficiency.value - levels_dropped)
                    if new_level < record.proficiency.value:
                        from_level_value = record.proficiency.value
                        record.proficiency = ProficiencyLevel(new_level)
                        record.assessment_history.append({
                            "timestamp": now,
                            "level": new_level,
                            "source": "decay",
                            "notes": f"Inactive for {idle_days:.0f} days",
                        })
                        await self._upsert_record(record)
                        decayed.append(record)
                        # AD-628a: SKILL_DECAY telemetry per decayed record
                        from probos.events import EventType as _ET
                        self._emit(_ET.SKILL_DECAY, {
                            "agent_id": record.agent_id,
                            "skill_id": record.skill_id,
                            "from_level": from_level_value,
                            "to_level": new_level,
                            "source": "decay",
                            "timestamp": now,
                            "reason": "idle_decay",
                        })
        return decayed

    async def get_profile(self, agent_id: str) -> SkillProfile:
        """Build the complete skill profile for an agent."""
        profile = SkillProfile(agent_id=agent_id)
        if not self._db:
            return profile

        async with self._db.execute(
            "SELECT * FROM agent_skills WHERE agent_id = ?", (agent_id,)
        ) as cur:
            async for row in cur:
                record = self._row_to_record(row)
                defn = self._registry.get_skill(record.skill_id) if self._registry else None
                if defn:
                    if defn.category == SkillCategory.PCC:
                        profile.pccs.append(record)
                    elif defn.category == SkillCategory.ROLE:
                        profile.role_skills.append(record)
                    else:
                        profile.acquired_skills.append(record)
                else:
                    profile.acquired_skills.append(record)
        return profile

    async def get_all_records(self, agent_id: str) -> list[AgentSkillRecord]:
        """Get all skill records for an agent."""
        if not self._db:
            return []
        records: list[AgentSkillRecord] = []
        async with self._db.execute(
            "SELECT * FROM agent_skills WHERE agent_id = ?", (agent_id,)
        ) as cur:
            async for row in cur:
                records.append(self._row_to_record(row))
        return records

    # ------------------------------------------------------------------
    # AD-428b v1: Development goals
    # ------------------------------------------------------------------

    async def add_development_goal(
        self,
        agent_id: str,
        skill_id: str,
        target_level: ProficiencyLevel,
        notes: str = "",
    ) -> dict[str, Any]:
        """AD-428b v1: set or replace a development goal for a (agent, skill).

        One goal per (agent_id, skill_id). Calling with the same skill_id
        replaces the existing goal. Returns the persisted goal as a dict.
        """
        now = time.time()
        if not self._db:
            return {
                "agent_id": agent_id,
                "skill_id": skill_id,
                "target_level": target_level.value,
                "set_at": now,
                "notes": notes,
            }
        await self._db.execute(
            "INSERT OR REPLACE INTO agent_development_goals "
            "(agent_id, skill_id, target_level, set_at, notes) "
            "VALUES (?, ?, ?, ?, ?)",
            (agent_id, skill_id, target_level.value, now, notes),
        )
        await self._db.commit()
        return {
            "agent_id": agent_id,
            "skill_id": skill_id,
            "target_level": target_level.value,
            "set_at": now,
            "notes": notes,
        }

    async def clear_development_goal(
        self, agent_id: str, skill_id: str
    ) -> bool:
        """AD-428b v1: remove a development goal. Returns True if a row was deleted."""
        if not self._db:
            return False
        cur = await self._db.execute(
            "DELETE FROM agent_development_goals WHERE agent_id = ? AND skill_id = ?",
            (agent_id, skill_id),
        )
        await self._db.commit()
        return cur.rowcount > 0  # type: ignore[attr-defined]

    async def get_development_goals(
        self, agent_id: str
    ) -> list[dict[str, Any]]:
        """AD-428b v1: return all development goals for an agent.

        Each entry: {skill_id, target_level (int 1-7), target_label, set_at, notes,
        current_level (int|None — agent's current proficiency on the skill)}.
        Sorted by skill_id.
        """
        if not self._db:
            return []
        # Build a map of agent's current proficiency per skill.
        records = await self.get_all_records(agent_id)
        current_by_skill = {r.skill_id: r.proficiency.value for r in records}
        goals: list[dict[str, Any]] = []
        async with self._db.execute(
            "SELECT skill_id, target_level, set_at, notes "
            "FROM agent_development_goals WHERE agent_id = ? ORDER BY skill_id",
            (agent_id,),
        ) as cur:
            async for row in cur:
                target_lvl = int(row["target_level"])
                try:
                    target_label = ProficiencyLevel(target_lvl).name
                except ValueError:
                    target_label = "UNKNOWN"
                goals.append({
                    "skill_id": row["skill_id"],
                    "target_level": target_lvl,
                    "target_label": target_label,
                    "set_at": row["set_at"],
                    "notes": row["notes"] or "",
                    "current_level": current_by_skill.get(row["skill_id"]),
                })
        return goals

    async def check_prerequisites(
        self, agent_id: str, skill_id: str,
    ) -> dict[str, Any]:
        """Check if an agent meets prerequisites for a skill.

        Returns {met: bool, missing: list[str]}.
        """
        if not self._registry:
            return {"met": True, "missing": []}
        defn = self._registry.get_skill(skill_id)
        if not defn or not defn.prerequisites:
            return {"met": True, "missing": []}
        missing: list[str] = []
        for prereq_id in defn.prerequisites:
            record = await self._get_record(agent_id, prereq_id)
            if not record or record.proficiency.value < ProficiencyLevel.APPLY.value:
                missing.append(prereq_id)
        return {"met": len(missing) == 0, "missing": missing}

    async def _get_record(self, agent_id: str, skill_id: str) -> AgentSkillRecord | None:
        if not self._db:
            return None
        async with self._db.execute(
            "SELECT * FROM agent_skills WHERE agent_id = ? AND skill_id = ?",
            (agent_id, skill_id),
        ) as cur:
            row = await cur.fetchone()
            if not row:
                return None
            return self._row_to_record(row)

    async def _upsert_record(self, record: AgentSkillRecord) -> None:
        if not self._db:
            return
        await self._db.execute(
            "INSERT OR REPLACE INTO agent_skills "
            "(agent_id, skill_id, proficiency, acquired_at, last_exercised, "
            "exercise_count, acquisition_source, suspended, assessment_history) "
            "VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?)",
            (record.agent_id, record.skill_id, record.proficiency.value,
             record.acquired_at, record.last_exercised, record.exercise_count,
             record.acquisition_source, int(record.suspended),
             json.dumps(record.assessment_history)),
        )
        await self._db.commit()

    # -------------------------------------------------------------------
    # Qualification tracking (AD-429b)
    # -------------------------------------------------------------------

    async def start_qualification(self, agent_id: str, path_id: str) -> QualificationRecord:
        """Start tracking a qualification path for an agent."""
        record = QualificationRecord(
            agent_id=agent_id,
            path_id=path_id,
            started_at=time.time(),
        )
        await self._save_qualification_record(record)
        return record

    async def evaluate_qualification(
        self, agent_id: str, path_id: str, ontology: Any
    ) -> QualificationRecord | None:
        """Evaluate an agent's current qualification status against path requirements."""
        if not ontology:
            return None
        parts = path_id.split("_to_")
        if len(parts) != 2:
            return None
        qual_path = ontology.get_qualification_path(parts[0], parts[1])
        if not qual_path:
            return None

        profile = await self.get_profile(agent_id)

        # Get role template for scope-based checks
        assignment = ontology.get_assignment_for_agent_by_id(agent_id)
        role_template = None
        if assignment:
            role_template = ontology.get_role_template(assignment.post_id)

        status: dict[str, bool] = {}
        for req in qual_path.requirements:
            key = f"{req.type}_{req.scope}"
            if req.scope == "all_pccs":
                pcc_records = profile.pccs
                if pcc_records:
                    status[key] = all(
                        r.proficiency.value >= req.min_proficiency for r in pcc_records
                    )
                else:
                    status[key] = False
            elif req.scope == "role_skills":
                role_records = profile.role_skills
                if req.min_count is not None:
                    count = sum(1 for r in role_records if r.proficiency.value >= req.min_proficiency)
                    status[key] = count >= req.min_count
                else:
                    status[key] = all(
                        r.proficiency.value >= req.min_proficiency for r in role_records
                    )
            elif req.scope == "required_role_skills":
                if role_template:
                    required_ids = {s.skill_id for s in role_template.required_skills}
                    required_records = [r for r in profile.role_skills if r.skill_id in required_ids]
                    status[key] = all(
                        r.proficiency.value >= req.min_proficiency for r in required_records
                    ) if required_records else False
                else:
                    status[key] = False

        # Get or create record
        record = await self.get_qualification_record(agent_id, path_id)
        if not record:
            record = QualificationRecord(
                agent_id=agent_id,
                path_id=path_id,
                started_at=time.time(),
            )
        record.requirement_status = status
        if record.is_complete() and not record.completed_at:
            record.completed_at = time.time()

        await self._save_qualification_record(record)
        return record

    async def get_qualification_record(self, agent_id: str, path_id: str) -> QualificationRecord | None:
        """Get a qualification record."""
        if not self._db:
            return None
        async with self._db.execute(
            "SELECT * FROM qualification_records WHERE agent_id = ? AND path_id = ?",
            (agent_id, path_id),
        ) as cur:
            row = await cur.fetchone()
            if not row:
                return None
            return QualificationRecord(
                agent_id=row["agent_id"],
                path_id=row["path_id"],
                started_at=row["started_at"],
                completed_at=row["completed_at"],
                requirement_status=json.loads(row["requirement_status"] or "{}"),
            )

    async def get_all_qualification_records(self, agent_id: str) -> list[QualificationRecord]:
        """Get all qualification records for an agent."""
        if not self._db:
            return []
        records: list[QualificationRecord] = []
        async with self._db.execute(
            "SELECT * FROM qualification_records WHERE agent_id = ?", (agent_id,)
        ) as cur:
            async for row in cur:
                records.append(QualificationRecord(
                    agent_id=row["agent_id"],
                    path_id=row["path_id"],
                    started_at=row["started_at"],
                    completed_at=row["completed_at"],
                    requirement_status=json.loads(row["requirement_status"] or "{}"),
                ))
        return records

    async def _save_qualification_record(self, record: QualificationRecord) -> None:
        """Persist a qualification record to SQLite."""
        if not self._db:
            return
        await self._db.execute(
            "INSERT OR REPLACE INTO qualification_records "
            "(agent_id, path_id, started_at, completed_at, requirement_status) "
            "VALUES (?, ?, ?, ?, ?)",
            (record.agent_id, record.path_id, record.started_at,
             record.completed_at, json.dumps(record.requirement_status)),
        )
        await self._db.commit()
