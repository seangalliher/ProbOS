"""AD-755 local template registry for recurring office tasks (OSS desktop scope)."""

from __future__ import annotations

import hashlib
import io
import json
import tempfile
from dataclasses import dataclass
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

from docx import Document
from openpyxl import load_workbook
from pptx import Presentation


@dataclass
class Template:
    """Template metadata and content bytes."""

    id: str
    name: str
    file_type: str
    content: bytes
    created_at: datetime
    used_count: int = 0


class TemplateRegistry:
    """Loads, lists, and instantiates local office templates."""

    def __init__(self, registry_dir: str) -> None:
        self._registry_dir = Path(registry_dir).expanduser()
        self._blobs_dir = self._registry_dir / "blobs"
        self._index_file = self._registry_dir / "index.json"
        self._registry_dir.mkdir(parents=True, exist_ok=True)
        self._blobs_dir.mkdir(parents=True, exist_ok=True)

    async def list_templates(self) -> list[Template]:
        """Return all locally registered templates."""
        entries = self._load_index_entries()
        templates: list[Template] = []
        for entry in entries:
            blob_path = self._blobs_dir / f"{entry['content_ref']}.bin"
            if not blob_path.exists():
                continue
            created_at = datetime.fromisoformat(entry["created_at"])
            templates.append(
                Template(
                    id=entry["id"],
                    name=entry["name"],
                    file_type=entry["file_type"],
                    content=blob_path.read_bytes(),
                    created_at=created_at,
                    used_count=int(entry.get("used_count", 0)),
                )
            )
        return templates

    async def create_from_template(self, template_id: str, **kwargs: Any) -> str:
        """Instantiate template with substitutions and return generated file path."""
        entries = self._load_index_entries()
        entry = next((item for item in entries if item.get("id") == template_id), None)
        if entry is None:
            raise ValueError(f"Template not found: {template_id}")

        blob_path = self._blobs_dir / f"{entry['content_ref']}.bin"
        if not blob_path.exists():
            raise FileNotFoundError(str(blob_path))

        substitutions = kwargs.get("substitutions") or {}
        output_path = kwargs.get("output_path")
        content = blob_path.read_bytes()
        file_type = str(entry["file_type"]).lower()

        generated = self._instantiate_bytes(content=content, file_type=file_type, substitutions=substitutions)

        if output_path:
            destination = Path(str(output_path))
            destination.parent.mkdir(parents=True, exist_ok=True)
            destination.write_bytes(generated)
        else:
            with tempfile.NamedTemporaryFile(suffix=f".{file_type}", delete=False) as tmp:
                destination = Path(tmp.name)
            destination.write_bytes(generated)

        entry["used_count"] = int(entry.get("used_count", 0)) + 1
        self._write_index_entries(entries)
        return str(destination)

    def _load_index_entries(self) -> list[dict[str, Any]]:
        if not self._index_file.exists():
            return []
        data = json.loads(self._index_file.read_text(encoding="utf-8"))
        if not isinstance(data, list):
            return []
        return [item for item in data if isinstance(item, dict)]

    def _write_index_entries(self, entries: list[dict[str, Any]]) -> None:
        self._index_file.write_text(
            json.dumps(entries, indent=2, sort_keys=True),
            encoding="utf-8",
        )

    def _instantiate_bytes(
        self,
        *,
        content: bytes,
        file_type: str,
        substitutions: dict[str, Any],
    ) -> bytes:
        if file_type == "docx":
            return self._instantiate_docx(content, substitutions)
        if file_type == "pptx":
            return self._instantiate_pptx(content, substitutions)
        if file_type == "xlsx":
            return self._instantiate_xlsx(content, substitutions)
        raise ValueError(f"Unsupported template file_type: {file_type}")

    def _instantiate_docx(self, content: bytes, substitutions: dict[str, Any]) -> bytes:
        doc = Document(io.BytesIO(content))
        for paragraph in doc.paragraphs:
            paragraph.text = self._apply_substitutions(paragraph.text, substitutions)
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    cell.text = self._apply_substitutions(cell.text, substitutions)

        buffer = io.BytesIO()
        doc.save(buffer)
        return buffer.getvalue()

    def _instantiate_pptx(self, content: bytes, substitutions: dict[str, Any]) -> bytes:
        deck = Presentation(io.BytesIO(content))
        for slide in deck.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    shape.text = self._apply_substitutions(shape.text, substitutions)

        buffer = io.BytesIO()
        deck.save(buffer)
        return buffer.getvalue()

    def _instantiate_xlsx(self, content: bytes, substitutions: dict[str, Any]) -> bytes:
        wb = load_workbook(io.BytesIO(content))
        for sheet in wb.worksheets:
            for row in sheet.iter_rows():
                for cell in row:
                    if isinstance(cell.value, str):
                        cell.value = self._apply_substitutions(cell.value, substitutions)

        buffer = io.BytesIO()
        wb.save(buffer)
        return buffer.getvalue()

    def _apply_substitutions(self, value: str, substitutions: dict[str, Any]) -> str:
        rendered = value
        for key, raw_value in substitutions.items():
            placeholder = "{{" + str(key) + "}}"
            rendered = rendered.replace(placeholder, str(raw_value))
        return rendered


def compute_template_ref(content: bytes) -> str:
    """Compute deterministic content ref for template blob storage."""
    return hashlib.sha256(content).hexdigest()


def serialize_template_entry(template: Template, content_ref: str) -> dict[str, Any]:
    """Serialize a template metadata row for index.json storage."""
    return {
        "id": template.id,
        "name": template.name,
        "file_type": template.file_type,
        "content_ref": content_ref,
        "created_at": template.created_at.astimezone(timezone.utc).isoformat(),
        "used_count": template.used_count,
    }
