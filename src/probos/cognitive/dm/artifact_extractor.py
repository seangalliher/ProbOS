"""AD-797 (Wave 197): artifact extractor for the DM reply pipeline.

Scans an agent's response text for *artifact-worthy* blocks — explicit
``<artifact>`` tags first, then large fenced-code blocks — persists the
bytes to ``AttachmentStore`` and metadata to ``ArtifactStore``, then
rewrites the response text with stub lines so the conversational
scrollback stays clean.

Two triggers:

1. **Explicit ``<artifact name="..." mime="...">...</artifact>`` tag.**
   Authoritative — used verbatim. Name is sanitized to
   ``[A-Za-z0-9._-]+`` (the name lands in a browser ``<a download>``
   attribute later, so strict validation is non-negotiable).

2. **Fenced code block** with ``>= fenced_threshold_lines`` lines
   (default 40, from ``CognitiveConfig.artifact_fenced_threshold_lines``).
   Mime + extension derived from the fence language tag via the
   12-language map below; default ``text/plain`` / ``.txt``. The block's
   first line may contain a ``# filename: helper.py`` comment (also
   ``// filename:`` and ``<!-- filename: ... -->``) to override the
   generated ``artifact-{N}.{ext}`` name.

Stub format (ASCII hyphen, not em-dash — UI regex matches the literal
hyphen)::

    [Artifact: helper.py v1 - 73 lines, text/x-python]
"""

from __future__ import annotations

import hashlib
import logging
import os
import re
from dataclasses import dataclass
from typing import Any

from probos.artifacts import Artifact, ArtifactStore

logger = logging.getLogger(__name__)


# Language tag → (mime, extension). Covers v1's 12 supported languages;
# additional mimes are AD-797e forward marker.
_LANG_MAP: dict[str, tuple[str, str]] = {
    "python": ("text/x-python", "py"),
    "py": ("text/x-python", "py"),
    "typescript": ("text/x-typescript", "ts"),
    "ts": ("text/x-typescript", "ts"),
    "javascript": ("text/x-javascript", "js"),
    "js": ("text/x-javascript", "js"),
    "markdown": ("text/markdown", "md"),
    "md": ("text/markdown", "md"),
    "bash": ("text/x-shellscript", "sh"),
    "sh": ("text/x-shellscript", "sh"),
    "json": ("application/json", "json"),
    "yaml": ("application/yaml", "yaml"),
    "yml": ("application/yaml", "yaml"),
    "html": ("text/html", "html"),
    "css": ("text/css", "css"),
    "sql": ("application/sql", "sql"),
    "rust": ("text/x-rust", "rs"),
    "rs": ("text/x-rust", "rs"),
    "go": ("text/x-go", "go"),
}

# Name sanitizer: filename chars only. Path separators stripped, other
# disallowed chars replaced with underscore. Empty result → caller skips.
_NAME_ALLOWED_RE = re.compile(r"[A-Za-z0-9._-]")
_NAME_DISALLOWED_RE = re.compile(r"[^A-Za-z0-9._-]")

# Filename hint inside the first line of a fenced block.
# Matches `# filename: foo.py`, `// filename: foo.ts`, `<!-- filename: foo.md -->`.
_FILENAME_COMMENT_RE = re.compile(
    r"""^\s*(?:
        \#\s*filename\s*:\s*(?P<py>[^\s]+)
      | //\s*filename\s*:\s*(?P<js>[^\s]+)
      | <!--\s*filename\s*:\s*(?P<html>[^\s]+)\s*-->
    )""",
    re.VERBOSE,
)

# Explicit <artifact name="..." mime="...">...</artifact>.
# Non-greedy body so two adjacent tags don't merge.
_ARTIFACT_TAG_RE = re.compile(
    r"<artifact\s+([^>]*?)>(.*?)</artifact>",
    re.DOTALL,
)
_ATTR_RE = re.compile(r"""(\w+)\s*=\s*(?:"([^"]*)"|'([^']*)')""")

# Fenced code block: ```lang\n ... \n```
_FENCE_RE = re.compile(
    r"```([A-Za-z0-9_+-]*)\s*\n(.*?)\n```",
    re.DOTALL,
)


@dataclass
class ExtractedArtifact:
    """One artifact extracted from response_text.

    ``source_span`` is the half-open ``(start, end)`` char-offset range
    of the *original* matched text in the unmodified ``response_text``.
    ``replace_with_stubs`` consumes these spans in reverse order so the
    earlier offsets stay valid as the text shrinks.
    """

    name: str
    mime: str
    content: bytes
    line_count: int
    source_span: tuple[int, int]


_DOCX_MIME = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
_PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
_XLSX_MIME = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"

_MD_HEADING_RE = re.compile(r"^(#{1,6})\s+(.*)$")
_MD_BULLET_RE = re.compile(r"^\s*[-*+\u2022]\s+(.*)$")
_MD_NUM_RE = re.compile(r"^\s*\d+[.)]\s+(.*)$")
_MD_BOLD_RE = re.compile(r"\*\*(.+?)\*\*")


def _add_md_runs(paragraph, text: str) -> None:
    """BF-647: split **bold** spans into runs so emphasis survives in docx."""
    pos = 0
    for m in _MD_BOLD_RE.finditer(text):
        if m.start() > pos:
            paragraph.add_run(text[pos:m.start()])
        paragraph.add_run(m.group(1)).bold = True
        pos = m.end()
    if pos < len(text):
        paragraph.add_run(text[pos:])


def _render_markdown_docx(doc, text: str) -> None:
    """BF-647: markdown-aware python-docx render (no external binary). Maps #
    headings, -/* bullets, numbered lists, and **bold** to real Word styles so
    quality is high without LibreOffice."""
    for raw in text.split("\n"):
        line = raw.rstrip()
        if not line.strip():
            continue
        h = _MD_HEADING_RE.match(line)
        if h:
            doc.add_heading(h.group(2).strip(), level=min(len(h.group(1)), 4))
            continue
        b = _MD_BULLET_RE.match(line)
        if b:
            _add_md_runs(doc.add_paragraph(style="List Bullet"), b.group(1).strip())
            continue
        n = _MD_NUM_RE.match(line)
        if n:
            _add_md_runs(doc.add_paragraph(style="List Number"), n.group(1).strip())
            continue
        _add_md_runs(doc.add_paragraph(), line.strip())


def _libreoffice_bin(explicit: str = "") -> str | None:
    """BF-646: resolve the soffice binary (explicit path > PATH). None if absent."""
    import shutil
    if explicit:
        return explicit if os.path.exists(explicit) else None
    return shutil.which("soffice") or shutil.which("soffice.exe")


def _to_docx_libreoffice(text: str, soffice: str) -> bytes | None:
    """BF-646: render text->docx via headless LibreOffice (higher fidelity).
    Wrap the prose as minimal HTML and `--convert-to docx`. None on any failure
    so the caller degrades to python-docx."""
    import subprocess
    import tempfile
    try:
        paras = "".join(f"<p>{p.strip()}</p>" for p in text.split("\n\n") if p.strip())
        html = f"<html><body>{paras or '<p></p>'}</body></html>"
        with tempfile.TemporaryDirectory() as d:
            src = os.path.join(d, "doc.html")
            with open(src, "w", encoding="utf-8") as f:
                f.write(html)
            subprocess.run(
                [soffice, "--headless", "--convert-to", "docx", "--outdir", d, src],
                check=True, capture_output=True, timeout=60,
            )
            out = os.path.join(d, "doc.docx")
            if os.path.exists(out):
                with open(out, "rb") as f:
                    return f.read()
    except Exception:
        logger.warning("BF-646: LibreOffice render failed; falling back", exc_info=True)
    return None


def _to_office_bytes(mime: str, content: bytes, backend: str = "python-docx", soffice_path: str = "") -> bytes:
    """BF-645/646/648: build a REAL Office binary from the agent's text. docx via
    python-docx (markdown-aware) or headless LibreOffice; pptx via python-pptx
    (a slide per ``#`` heading); xlsx via openpyxl (CSV/tab rows). All honest-
    degrade to the original bytes on failure."""
    text = content.decode("utf-8", errors="replace")
    if mime == _DOCX_MIME:
        if backend == "libreoffice":
            soffice = _libreoffice_bin(soffice_path)
            if soffice:
                rendered = _to_docx_libreoffice(text, soffice)
                if rendered:
                    return rendered
            logger.warning("BF-646: libreoffice backend unavailable; using python-docx")
        try:
            from io import BytesIO
            from docx import Document
            doc = Document()
            _render_markdown_docx(doc, text)
            buf = BytesIO()
            doc.save(buf)
            return buf.getvalue()
        except Exception:
            logger.warning("BF-645: docx render failed; persisting raw text", exc_info=True)
            return content
    if mime == _PPTX_MIME:
        return _to_pptx(text) or content
    if mime == _XLSX_MIME:
        return _to_xlsx(text) or content
    return content


def _to_pptx(text: str) -> bytes | None:
    """BF-648: a deck where each ``#``/``##`` heading starts a slide; bullets and
    prose become the body. None on failure -> caller keeps raw bytes."""
    try:
        from io import BytesIO
        from pptx import Presentation
        from pptx.util import Inches
        prs = Presentation()
        blank = prs.slide_layouts[1]  # title + content
        title, body = "", []

        def _flush():
            if not title and not body:
                return
            s = prs.slides.add_slide(blank)
            s.shapes.title.text = title or "Slide"
            tf = s.placeholders[1].text_frame
            tf.text = body[0] if body else ""
            for line in body[1:]:
                tf.add_paragraph().text = line

        for raw in text.split("\n"):
            line = raw.rstrip()
            h = _MD_HEADING_RE.match(line)
            if h:
                _flush()
                title, body = h.group(2).strip(), []
            elif line.strip():
                b = _MD_BULLET_RE.match(line) or _MD_NUM_RE.match(line)
                body.append((b.group(1) if b else line).strip())
        _flush()
        if not prs.slides:
            prs.slides.add_slide(blank).shapes.title.text = text[:80] or "Slide"
        buf = BytesIO()
        prs.save(buf)
        return buf.getvalue()
    except Exception:
        logger.warning("BF-648: pptx render failed", exc_info=True)
        return None


def _to_xlsx(text: str) -> bytes | None:
    """BF-648: rows from comma/tab-separated lines. None on failure."""
    try:
        from io import BytesIO
        from openpyxl import Workbook
        wb = Workbook()
        ws = wb.active
        for raw in text.split("\n"):
            line = raw.rstrip()
            if not line.strip():
                continue
            cells = line.split("\t") if "\t" in line else line.split(",")
            ws.append([c.strip() for c in cells])
        buf = BytesIO()
        wb.save(buf)
        return buf.getvalue()
    except Exception:
        logger.warning("BF-648: xlsx render failed", exc_info=True)
        return None


def _sanitize_name(raw: str) -> str:
    """Strip path separators; replace disallowed chars with ``_``.

    Returns ``""`` when the result has no allowed chars at all
    (caller logs + skips). The name will be used in a browser
    ``<a download={name}>`` attribute — strict validation is required
    for safety.
    """
    if not raw:
        return ""
    # Strip directory traversal components first.
    just_name = raw.replace("\\", "/").rsplit("/", 1)[-1]
    sanitized = _NAME_DISALLOWED_RE.sub("_", just_name).strip(".")
    if not _NAME_ALLOWED_RE.search(sanitized):
        return ""
    return sanitized


def _parse_tag_attrs(attr_blob: str) -> dict[str, str]:
    attrs: dict[str, str] = {}
    for m in _ATTR_RE.finditer(attr_blob):
        key = m.group(1)
        val = m.group(2) if m.group(2) is not None else m.group(3) or ""
        attrs[key] = val
    return attrs


def _line_count(blob: bytes, mime: str) -> int:
    if mime.startswith("image/") or mime == "text/uri-list":
        return 0
    try:
        text = blob.decode("utf-8", errors="replace")
    except Exception:
        return 0
    if not text:
        return 0
    # Count newlines + 1 if the final line has no trailing \n.
    n = text.count("\n")
    if not text.endswith("\n"):
        n += 1
    return n


def _filename_hint(first_line: str) -> str | None:
    m = _FILENAME_COMMENT_RE.match(first_line)
    if not m:
        return None
    return m.group("py") or m.group("js") or m.group("html")


def has_explicit_artifact_marker(response_text: str) -> bool:
    """Whether the reply carries an explicit ``<artifact>`` tag (pass 1).

    AD-1285 (#1087): separates a save the agent deliberately asked for from the
    pass-2 lift of any long fenced block, which it never claimed to save. Only
    the former may be reported as a write that was attempted and failed --
    treating a passive lift the same way would let the text-blind write-claim
    guard append a save disclosure to a reply describing no save.
    """
    return bool(_ARTIFACT_TAG_RE.search(response_text or ""))


def extract_artifacts(
    response_text: str,
    *,
    fenced_threshold_lines: int = 40,
    existing_unnamed_count: int = 0,
) -> list[ExtractedArtifact]:
    """Two-pass scan of ``response_text``.

    Pass 1 (explicit ``<artifact>`` tags): authoritative. Each tag with a
    valid (sanitizable) name and a mime attribute produces one
    ``ExtractedArtifact``.

    Pass 2 (fenced code blocks): scans the *remaining* body (with tag
    spans masked out) for ``` blocks with ``>= fenced_threshold_lines``
    lines. Name resolution: filename-comment hint → otherwise
    ``artifact-{N}.{ext}`` where ``N = existing_unnamed_count + 1`` and
    increments per unnamed extraction in this call.

    Returns artifacts in source-position order so the stub replacement
    pass produces stable output.
    """
    if not response_text:
        return []

    extracted: list[ExtractedArtifact] = []

    # Track masked spans so the fenced-pass skips characters inside tags.
    masked_spans: list[tuple[int, int]] = []

    # --- Pass 1: explicit <artifact> tags ---
    for m in _ARTIFACT_TAG_RE.finditer(response_text):
        attrs = _parse_tag_attrs(m.group(1))
        raw_name = attrs.get("name", "")
        mime = attrs.get("mime", "").strip()
        if not mime:
            logger.warning(
                "AD-797: <artifact> tag missing mime attribute; skipping (span=%s)",
                m.span(),
            )
            masked_spans.append(m.span())
            continue
        name = _sanitize_name(raw_name)
        if not name:
            logger.warning(
                "AD-797: <artifact> tag has unsafe/empty name %r; skipping",
                raw_name,
            )
            masked_spans.append(m.span())
            continue
        body = m.group(2)
        # For text mimes, strip a single leading + trailing newline so
        # the body matches what the agent wrote between the tags.
        if not mime.startswith("image/"):
            if body.startswith("\n"):
                body = body[1:]
            if body.endswith("\n"):
                body = body[:-1]
        blob = body.encode("utf-8")
        extracted.append(
            ExtractedArtifact(
                name=name,
                mime=mime,
                content=blob,
                line_count=_line_count(blob, mime),
                source_span=m.span(),
            )
        )
        masked_spans.append(m.span())

    # --- Pass 2: fenced code blocks not inside an <artifact> tag ---
    def _is_masked(span: tuple[int, int]) -> bool:
        s, e = span
        for ms, me in masked_spans:
            if not (e <= ms or s >= me):
                return True
        return False

    unnamed_n = existing_unnamed_count
    for m in _FENCE_RE.finditer(response_text):
        if _is_masked(m.span()):
            continue
        lang = (m.group(1) or "").strip().lower()
        body = m.group(2)
        line_count = body.count("\n") + (0 if body.endswith("\n") else 1)
        if line_count < fenced_threshold_lines:
            continue
        mime, ext = _LANG_MAP.get(lang, ("text/plain", "txt"))
        # Filename hint on first line of the block (also: respect ``ext``
        # inferred from the hint over the lang map).
        first_line = body.split("\n", 1)[0] if body else ""
        hint = _filename_hint(first_line)
        if hint:
            sanitized = _sanitize_name(hint)
            if sanitized:
                name = sanitized
                # Re-derive mime from hint extension when lang map disagrees.
                if "." in sanitized:
                    hint_ext = sanitized.rsplit(".", 1)[-1].lower()
                    for k, (kmime, kext) in _LANG_MAP.items():
                        if kext == hint_ext:
                            mime = kmime
                            break
            else:
                unnamed_n += 1
                name = f"artifact-{unnamed_n}.{ext}"
        else:
            unnamed_n += 1
            name = f"artifact-{unnamed_n}.{ext}"
        blob = body.encode("utf-8")
        extracted.append(
            ExtractedArtifact(
                name=name,
                mime=mime,
                content=blob,
                line_count=_line_count(blob, mime),
                source_span=m.span(),
            )
        )

    extracted.sort(key=lambda e: e.source_span[0])
    return extracted


async def replace_with_stubs(
    response_text: str,
    extracted: list[ExtractedArtifact],
    *,
    artifact_store: ArtifactStore,
    attachment_store: Any,
    thread_id: str,
    created_by: str,
    office_backend: str = "python-docx",
    libreoffice_path: str = "",
) -> tuple[str, list[Artifact]]:
    """Persist each extracted artifact and rewrite ``response_text``.

    For each extracted block, in source order:

    1. Compute ``content_hash = sha256(blob).hexdigest()``.
    2. ``await attachment_store.write(hash, blob, mime,
       origin="agent_artifact")`` (idempotent).
    3. ``artifact_store.add_version(...)`` — auto-assigns version 1,
       2, 3, ... per ``(thread_id, name)``.
    4. Replace the source span with the stub line.

    Replacements are applied in reverse order so earlier offsets stay
    valid as the text shrinks. Returns the rewritten text + the
    persisted ``Artifact`` rows in source order.
    """
    if not extracted:
        return response_text, []

    persisted: list[Artifact] = []
    # First persist; collect (span, stub) pairs.
    stubs: list[tuple[tuple[int, int], str]] = []
    for ex in extracted:
        # BF-645: render docx prose to a real Word binary before hashing/persist.
        blob = _to_office_bytes(ex.mime, ex.content, office_backend, libreoffice_path)
        content_hash = hashlib.sha256(blob).hexdigest()
        try:
            await attachment_store.write(
                content_hash,
                blob,
                ex.mime,
                origin="agent_artifact",
            )
        except Exception as exc:
            logger.warning(
                "AD-797: AttachmentStore.write failed for %s (hash=%s); "
                "skipping artifact persist",
                ex.name, content_hash[:12], exc_info=True,
            )
            continue
        try:
            artifact = artifact_store.add_version(
                thread_id=thread_id,
                name=ex.name,
                content_hash=content_hash,
                mime=ex.mime,
                size_bytes=len(blob),
                created_by=created_by,
            )
        except Exception:
            logger.warning(
                "AD-797: ArtifactStore.add_version failed for %s; "
                "skipping stub replacement",
                ex.name, exc_info=True,
            )
            continue
        persisted.append(artifact)
        stub = (
            f"[Artifact: {artifact.name} v{artifact.version} - "
            f"{ex.line_count} lines, {artifact.mime}]"
        )
        stubs.append((ex.source_span, stub))

    # Apply stubs in reverse source order so earlier offsets stay valid.
    stubs.sort(key=lambda s: s[0][0], reverse=True)
    new_text = response_text
    for (start, end), stub in stubs:
        new_text = new_text[:start] + stub + new_text[end:]

    return new_text, persisted
